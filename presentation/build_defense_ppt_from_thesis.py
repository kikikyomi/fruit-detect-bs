from __future__ import annotations

import shutil
from pathlib import Path

from PIL import Image, ImageOps, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PPT_PATH = Path(r"D:\毕设\fruit-veg-dete\presentation\fruit_veg_defense.pptx")
BACKUP_PATH = PPT_PATH.with_name("fruit_veg_defense.before_thesis_rebuild.pptx")
ASSET_DIR = PPT_PATH.parent / "_thesis_assets"

DETECTION_SAMPLE = Path(r"D:\毕设\fruit-veg-dete\presentation\_ppt_media\ppt\media\image1.png")
TRACK_SAMPLE = Path(r"D:\毕设\fruit-veg-dete\presentation\_ppt_media\ppt\media\image2.png")
YOLO_SLIDE = Path(r"D:\毕设\fruit-veg-dete\presentation\_ppt_media\ppt\media\image3.png")
DEEPSORT_SLIDE = Path(r"D:\毕设\fruit-veg-dete\presentation\_ppt_media\ppt\media\image4.png")

DATA_STATS_IMAGE = Path(r"D:\毕设\fruit-veg-dete\presentation\_docx_media\word\media\image3.png")
PREPROCESS_IMAGE = Path(r"D:\毕设\fruit-veg-dete\presentation\_docx_media\word\media\image4.png")
RESULTS_IMAGE = Path(r"D:\毕设\fruit-veg-dete\presentation\_docx_media\word\media\image10.png")
BUSINESS_FLOW_IMAGE = Path(r"D:\毕设\fruit-veg-dete\presentation\_docx_media\word\media\image11.png")
SYSTEM_FLOW_IMAGE = Path(r"D:\毕设\fruit-veg-dete\presentation\_docx_media\word\media\image12.png")
UI_IMAGE = Path(r"D:\毕设\fruit-veg-dete\presentation\_docx_media\word\media\image13.png")


COLORS = {
    "navy": RGBColor(27, 44, 84),
    "blue": RGBColor(49, 107, 184),
    "green": RGBColor(84, 155, 98),
    "light_green": RGBColor(228, 242, 230),
    "orange": RGBColor(228, 134, 73),
    "light_orange": RGBColor(250, 235, 223),
    "bg": RGBColor(247, 249, 252),
    "text": RGBColor(46, 57, 75),
    "muted": RGBColor(108, 118, 136),
    "line": RGBColor(217, 223, 233),
    "white": RGBColor(255, 255, 255),
    "soft_blue": RGBColor(235, 242, 252),
    "soft_red": RGBColor(252, 235, 232),
}


def ensure_exists(path: Path) -> None:
    if not path.exists():
        raise FileNotFoundError(path)


def set_paragraph_style(paragraph, *, size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Microsoft YaHei"
        run.font.color.rgb = color or COLORS["text"]


def add_base_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.32)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS["navy"]
    top_bar.line.fill.background()
    return slide


def add_title(slide, number: int, title: str) -> None:
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(0.52),
        Inches(0.92),
        Inches(0.46),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLORS["orange"]
    badge.line.fill.background()
    badge.text_frame.text = f"{number:02d}"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_paragraph_style(
        badge.text_frame.paragraphs[0], size=20, bold=True, color=COLORS["white"]
    )

    title_box = slide.shapes.add_textbox(
        Inches(1.55), Inches(0.48), Inches(8.8), Inches(0.62)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    set_paragraph_style(p, size=26, bold=True, color=COLORS["navy"])


def add_section_chip(slide, text: str) -> None:
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(0.68),
        Inches(1.5),
        Inches(0.42),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = COLORS["light_orange"]
    chip.line.fill.background()
    chip.text_frame.text = text
    chip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_paragraph_style(chip.text_frame.paragraphs[0], size=12, bold=True, color=COLORS["navy"])


def rounded_card(slide, left: float, top: float, width: float, height: float, fill: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = COLORS["line"]
    return shape


def add_bullet_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list[str],
    *,
    fill: RGBColor = COLORS["white"],
    title_size: int = 18,
    text_size: int = 15,
) -> None:
    card = rounded_card(slide, left, top, width, height, fill)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)

    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_style(p, size=title_size, bold=True, color=COLORS["navy"])

    for bullet in bullets:
        para = tf.add_paragraph()
        para.text = f"• {bullet}"
        para.space_after = Pt(6)
        set_paragraph_style(para, size=text_size, color=COLORS["text"])


def add_stat_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    value: str,
    label: str,
    *,
    fill: RGBColor = COLORS["white"],
) -> None:
    card = rounded_card(slide, left, top, width, height, fill)
    tf = card.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.text = value
    set_paragraph_style(p1, size=20, bold=True, color=COLORS["navy"])

    p2 = tf.add_paragraph()
    p2.text = label
    set_paragraph_style(p2, size=10, color=COLORS["muted"])


def add_center_text_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    subtitle: str | None = None,
    *,
    fill: RGBColor = COLORS["white"],
    title_size: int = 16,
) -> None:
    card = rounded_card(slide, left, top, width, height, fill)
    tf = card.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    set_paragraph_style(p1, size=title_size, bold=True, color=COLORS["navy"])

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.alignment = PP_ALIGN.CENTER
        set_paragraph_style(p2, size=10, color=COLORS["muted"])


def fit_image(source: Path, output: Path, width_px: int, height_px: int) -> Path:
    with Image.open(source) as image:
        image = ImageOps.exif_transpose(image).convert("RGB")
        image.thumbnail((width_px, height_px), Image.LANCZOS)
        canvas = Image.new("RGB", (width_px, height_px), "white")
        x = (width_px - image.width) // 2
        y = (height_px - image.height) // 2
        canvas.paste(image, (x, y))
        canvas.save(output, quality=95)
    return output


def build_cover_collage(output: Path) -> Path:
    canvas = Image.new("RGB", (1100, 1100), "white")
    draw = ImageDraw.Draw(canvas)
    draw.rounded_rectangle((18, 18, 1082, 1082), radius=36, fill=(255, 255, 255))

    with Image.open(DETECTION_SAMPLE) as det:
        det = ImageOps.exif_transpose(det).convert("RGB")
        det.thumbnail((1030, 610), Image.LANCZOS)
        dx = (1100 - det.width) // 2
        canvas.paste(det, (dx, 38))

    with Image.open(UI_IMAGE) as ui:
        ui = ImageOps.exif_transpose(ui).convert("RGB")
        ui.thumbnail((1030, 390), Image.LANCZOS)
        ux = (1100 - ui.width) // 2
        canvas.paste(ui, (ux, 660))

    draw.rounded_rectangle((38, 38, 1062, 640), radius=24, outline=(217, 223, 233), width=2)
    draw.rounded_rectangle((38, 660, 1062, 1048), radius=24, outline=(217, 223, 233), width=2)
    canvas.save(output, quality=95)
    return output


def add_picture_card(
    slide,
    image_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    caption: str | None = None,
) -> None:
    card = rounded_card(slide, left, top, width, height, COLORS["white"])
    inner_left = left + 0.08
    inner_top = top + 0.08
    inner_width = width - 0.16
    inner_height = height - (0.28 if caption else 0.16)
    slide.shapes.add_picture(
        str(image_path),
        Inches(inner_left),
        Inches(inner_top),
        Inches(inner_width),
        Inches(inner_height),
    )
    if caption:
        box = slide.shapes.add_textbox(
            Inches(left + 0.12), Inches(top + height - 0.26), Inches(width - 0.24), Inches(0.15)
        )
        p = box.text_frame.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_style(p, size=9, color=COLORS["muted"])


def add_bottom_note(slide, text: str, *, fill: RGBColor = COLORS["soft_blue"]) -> None:
    card = rounded_card(slide, 0.55, 6.25, 12.2, 0.52, fill)
    tf = card.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_style(p, size=11, color=COLORS["muted"])


def build_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(exist_ok=True)
    ensure_exists(DETECTION_SAMPLE)
    ensure_exists(TRACK_SAMPLE)
    ensure_exists(YOLO_SLIDE)
    ensure_exists(DEEPSORT_SLIDE)
    ensure_exists(DATA_STATS_IMAGE)
    ensure_exists(PREPROCESS_IMAGE)
    ensure_exists(RESULTS_IMAGE)
    ensure_exists(BUSINESS_FLOW_IMAGE)
    ensure_exists(SYSTEM_FLOW_IMAGE)
    ensure_exists(UI_IMAGE)

    assets = {
        "cover": build_cover_collage(ASSET_DIR / "cover_collage.png"),
        "dataset": fit_image(DATA_STATS_IMAGE, ASSET_DIR / "dataset_stats.png", 1280, 760),
        "preprocess": fit_image(PREPROCESS_IMAGE, ASSET_DIR / "preprocess.png", 920, 760),
        "results": fit_image(RESULTS_IMAGE, ASSET_DIR / "results_compare.png", 1340, 820),
        "system": fit_image(SYSTEM_FLOW_IMAGE, ASSET_DIR / "system_flow.png", 1500, 820),
        "business": fit_image(BUSINESS_FLOW_IMAGE, ASSET_DIR / "business_flow.png", 1180, 760),
        "ui": fit_image(UI_IMAGE, ASSET_DIR / "ui.png", 1180, 760),
        "yolo": fit_image(YOLO_SLIDE, ASSET_DIR / "yolo_slide.png", 1180, 760),
        "deepsort": fit_image(DEEPSORT_SLIDE, ASSET_DIR / "deepsort_slide.png", 1180, 760),
        "image_result": fit_image(DETECTION_SAMPLE, ASSET_DIR / "image_result.png", 1400, 880),
        "video_result": fit_image(TRACK_SAMPLE, ASSET_DIR / "video_result.png", 1040, 860),
    }
    return assets


def build_presentation(assets: dict[str, Path]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = add_base_slide(prs)
    add_section_chip(slide, "毕业设计答辩")
    hero = rounded_card(slide, 7.0, 0.95, 5.65, 5.55, COLORS["white"])
    slide.shapes.add_picture(str(assets["cover"]), Inches(7.08), Inches(1.03), Inches(5.49), Inches(5.39))

    title_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.28), Inches(5.95), Inches(1.7))
    title_box.text_frame.word_wrap = True
    p1 = title_box.text_frame.paragraphs[0]
    p1.text = "面向果蔬智能采摘的动态目标定位与跟踪系统设计与实现"
    set_paragraph_style(p1, size=28, bold=True, color=COLORS["navy"])
    p2 = title_box.text_frame.add_paragraph()
    p2.text = "Design and Implementation of a Dynamic Target Localization and Tracking System for Intelligent Fruit and Vegetable Harvesting"
    set_paragraph_style(p2, size=12, color=COLORS["muted"])

    info = rounded_card(slide, 0.62, 3.18, 5.35, 2.05, COLORS["white"])
    tf = info.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.12)
    for idx, text in enumerate(
        [
            "姓名：张帅",
            "学号：220224337",
            "导师：赵昀杰",
            "日期：2026年04月22日",
        ]
    ):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = text
        para.space_after = Pt(8)
        set_paragraph_style(para, size=20, color=COLORS["text"])

    sub = rounded_card(slide, 0.62, 5.55, 5.35, 0.72, COLORS["soft_blue"])
    p = sub.text_frame.paragraphs[0]
    p.text = "围绕目标检测、跨帧跟踪和系统部署三个层面展开答辩"
    p.alignment = PP_ALIGN.CENTER
    sub.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_paragraph_style(p, size=12, bold=True, color=COLORS["navy"])

    # Slide 2
    slide = add_base_slide(prs)
    add_title(slide, 2, "目录")
    sections = [
        ("01", "研究背景与研究目标", COLORS["soft_blue"]),
        ("02", "数据集与模型选型", COLORS["light_orange"]),
        ("03", "系统设计与实现", COLORS["light_green"]),
        ("04", "实验结果与结论", COLORS["soft_red"]),
    ]
    positions = [(0.9, 1.85), (6.8, 1.85), (0.9, 3.7), (6.8, 3.7)]
    for (num, text, fill), (left, top) in zip(sections, positions):
        card = rounded_card(slide, left, top, 5.35, 1.35, fill)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.text = num
        set_paragraph_style(p1, size=13, bold=True, color=COLORS["orange"])
        p2 = tf.add_paragraph()
        p2.text = text
        set_paragraph_style(p2, size=22, bold=True, color=COLORS["navy"])
    add_bottom_note(slide, "整体汇报控制在 10 分钟左右，重点突出模型选型依据、系统落地和实验效果。")

    # Slide 3
    slide = add_base_slide(prs)
    add_title(slide, 3, "研究背景与研究目标")
    add_bullet_card(
        slide,
        0.7,
        1.45,
        4.45,
        3.15,
        "问题背景",
        [
            "采摘场景中目标密集、尺度变化大，且存在枝叶遮挡与复杂光照。",
            "仅做单帧检测难以支撑连续定位、计数和轨迹分析。",
            "论文目标是构建“检测 + 跟踪 + 可视化 + 存储”一体化系统。",
            "最终支持图片、视频、摄像头三类输入方式。",
        ],
        text_size=15,
    )
    add_center_text_card(slide, 5.45, 1.55, 2.15, 1.35, "研究对象", "13 类果蔬", fill=COLORS["soft_blue"])
    add_center_text_card(slide, 7.9, 1.55, 2.15, 1.35, "核心方法", "YOLO + DeepSORT", fill=COLORS["light_orange"])
    add_center_text_card(slide, 10.35, 1.55, 2.15, 1.35, "工程目标", "Web 系统部署", fill=COLORS["light_green"])
    add_bullet_card(
        slide,
        5.45,
        3.1,
        6.95,
        2.25,
        "本文主要工作",
        [
            "完成 13 类果蔬数据清洗、划分和增强处理。",
            "对 YOLOv8n、YOLOv8s、YOLOv8m、YOLO11s 进行训练比较并完成主模型选型。",
            "融合 DeepSORT 与 IoU 回退策略，实现跨帧关联与轨迹可视化。",
            "基于 Vue3、FastAPI 和 SQLite 实现可部署的检测跟踪系统。",
        ],
        fill=COLORS["white"],
        title_size=17,
        text_size=13,
    )
    add_bottom_note(slide, "答辩时可先讲“为什么要做”，再讲“做了什么”和“最终达成了什么”。", fill=COLORS["light_green"])

    # Slide 4
    slide = add_base_slide(prs)
    add_title(slide, 4, "数据集与预处理设计")
    add_stat_card(slide, 0.78, 1.45, 2.55, 0.95, "21428", "图像数量", fill=COLORS["soft_blue"])
    add_stat_card(slide, 3.55, 1.45, 2.55, 0.95, "13", "目标类别数", fill=COLORS["light_orange"])
    add_stat_card(slide, 6.32, 1.45, 2.55, 0.95, "252452", "标注框数量", fill=COLORS["light_green"])
    add_stat_card(slide, 9.09, 1.45, 2.95, 0.95, "95%≤32", "每图框数上界", fill=COLORS["soft_blue"])
    add_picture_card(slide, assets["dataset"], 0.72, 2.65, 6.95, 3.45, caption="数据集统计图：样本规模、类别分布与框密度")
    add_picture_card(slide, assets["preprocess"], 7.85, 2.65, 4.95, 3.45, caption="预处理流程：清洗、划分与 data.yaml 构建")
    add_bottom_note(
        slide,
        "标签统一采用 YOLO 格式，按 8:1:1 划分训练/验证/测试集；训练阶段启用 Mosaic、MixUp、随机缩放和 HSV 扰动。",
    )

    # Slide 5
    slide = add_base_slide(prs)
    add_title(slide, 5, "模型对比与主模型选型")
    add_picture_card(slide, assets["results"], 0.72, 2.02, 7.7, 4.2, caption="四种模型训练结果曲线对比")
    add_stat_card(slide, 8.7, 2.05, 1.72, 0.95, "0.771", "YOLOv8n mAP@0.5", fill=COLORS["soft_blue"])
    add_stat_card(slide, 10.56, 2.05, 1.72, 0.95, "0.805", "YOLOv8s mAP@0.5", fill=COLORS["light_orange"])
    add_stat_card(slide, 8.7, 3.15, 1.72, 0.95, "0.829", "YOLOv8m mAP@0.5", fill=COLORS["light_green"])
    add_stat_card(slide, 10.56, 3.15, 1.72, 0.95, "0.820", "YOLO11s mAP@0.5", fill=COLORS["soft_blue"])
    add_bullet_card(
        slide,
        8.7,
        4.35,
        3.58,
        1.88,
        "选型结论",
        [
            "YOLOv8m 后期收敛更平滑，整体鲁棒性更好。",
            "PR 曲线和混淆矩阵表现均较稳定。",
            "因此系统部署阶段选择 YOLOv8m 作为主检测模型。",
        ],
        fill=COLORS["white"],
        title_size=17,
        text_size=12,
    )
    add_bottom_note(slide, "训练配置：输入尺寸 640×640，batch=16，epoch=150，优化器采用 AdamW，并启用混合精度训练。")

    # Slide 6
    slide = add_base_slide(prs)
    add_title(slide, 6, "系统总体架构与业务流程")
    add_picture_card(slide, assets["system"], 0.72, 1.52, 12.0, 3.95, caption="系统总体业务流程：核心识别功能、结果输出与后期管理")
    add_stat_card(slide, 0.9, 5.72, 3.65, 0.95, "Vue3 + TypeScript", "前端交互与可视化", fill=COLORS["soft_blue"])
    add_stat_card(slide, 4.83, 5.72, 3.65, 0.95, "FastAPI + YOLO + DeepSORT", "后端推理与跟踪服务", fill=COLORS["light_orange"])
    add_stat_card(slide, 8.76, 5.72, 3.65, 0.95, "SQLite", "结果持久化与历史记录管理", fill=COLORS["light_green"])

    # Slide 7
    slide = add_base_slide(prs)
    add_title(slide, 7, "核心功能模块与前端界面")
    add_picture_card(slide, assets["business"], 0.72, 1.5, 6.0, 4.75, caption="功能流程：图片检测、视频跟踪、摄像头识别与历史数据管理")
    add_picture_card(slide, assets["ui"], 6.95, 1.5, 5.8, 4.75, caption="前端界面：上传、阈值调节、画布展示与检测列表")
    add_bottom_note(
        slide,
        "系统支持图片拖拽上传、视频文件提交、摄像头调用与历史记录回放，目标是降低非专业用户的操作门槛。",
    )

    # Slide 8
    slide = add_base_slide(prs)
    add_title(slide, 8, "模型与跟踪策略设计")
    add_picture_card(slide, assets["yolo"], 0.72, 1.48, 6.0, 3.9, caption="检测端：YOLO 模型结构解析")
    add_picture_card(slide, assets["deepsort"], 6.95, 1.48, 5.8, 3.9, caption="跟踪端：DeepSORT 多目标跟踪机制")
    add_bullet_card(
        slide,
        0.72,
        5.55,
        12.0,
        1.15,
        "优化策略",
        [
            "检测端比较四种 YOLO 模型后选择 YOLOv8m；跟踪端由 DeepSORT 维护 TrackID，并加入 IoU 回退策略应对遮挡与匹配失败。",
            "关键参数设置为 max_age=30、n_init=3、max_iou_distance=0.7、max_cosine_distance=0.2，轨迹缓存长度为 64。",
        ],
        fill=COLORS["white"],
        title_size=16,
        text_size=12,
    )

    # Slide 9
    slide = add_base_slide(prs)
    add_title(slide, 9, "图片识别模块与样例结果")
    add_bullet_card(
        slide,
        0.72,
        1.55,
        4.12,
        3.18,
        "模块说明",
        [
            "上传单张图片后，调用 YOLOv8m 完成果蔬定位与分类。",
            "页面同步展示类别、置信度和结构化检测列表。",
            "结果图与检测记录写入 SQLite，支持后续查询与复核。",
            "当前展示样例共识别 6 个目标。",
        ],
        text_size=14,
    )
    add_stat_card(slide, 0.72, 4.95, 1.95, 0.88, "6", "当前样例识别框", fill=COLORS["soft_blue"])
    add_stat_card(slide, 2.85, 4.95, 1.99, 0.88, "标签 + 置信度", "结果叠加可视化", fill=COLORS["light_orange"])
    add_picture_card(slide, assets["image_result"], 5.12, 1.42, 7.62, 4.72, caption="图片识别样例：检测框、类别和置信度同步输出")
    add_bottom_note(slide, "图片检测模块重点体现系统的可解释性：不仅有最终结果图，还保留了结构化识别数据。")

    # Slide 10
    slide = add_base_slide(prs)
    add_title(slide, 10, "视频跟踪模块与样例结果")
    add_picture_card(slide, assets["video_result"], 0.72, 1.45, 5.3, 4.7, caption="视频跟踪样例：跨帧 ID 关联与轨迹展示")
    add_stat_card(slide, 6.3, 1.58, 2.65, 0.98, "720", "样例视频总帧数", fill=COLORS["soft_blue"])
    add_stat_card(slide, 9.05, 1.58, 2.65, 0.98, "360", "抽样处理帧数", fill=COLORS["light_orange"])
    add_stat_card(slide, 6.3, 2.78, 2.65, 0.98, "3468", "检测次数", fill=COLORS["light_green"])
    add_stat_card(slide, 9.05, 2.78, 2.65, 0.98, "15.17 FPS", "样例处理速度", fill=COLORS["soft_blue"])
    add_bullet_card(
        slide,
        6.3,
        4.05,
        5.4,
        1.7,
        "模块说明",
        [
            "视频先完成解码与抽帧，再将检测结果送入 DeepSORT 做 TrackID 关联。",
            "系统输出轨迹视频、关键帧与统计摘要，平均视频处理速度约为 9.61 FPS。",
            "IoU 回退策略降低了短时遮挡和抖动场景下的 ID 中断风险。",
        ],
        title_size=16,
        text_size=12,
    )
    add_bottom_note(slide, "视频模块体现了“检测 + 跟踪”闭环，是后续果实计数和运动轨迹分析的基础。")

    # Slide 11
    slide = add_base_slide(prs)
    add_title(slide, 11, "系统测试、结论与展望")
    banner = rounded_card(slide, 0.7, 1.38, 12.0, 0.92, COLORS["soft_blue"])
    banner.text_frame.text = "系统完成了从数据预处理、模型训练到前后端部署的完整闭环，可较好满足果蔬采摘场景下的基础视觉感知需求。"
    banner.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_paragraph_style(banner.text_frame.paragraphs[0], size=18, bold=True, color=COLORS["navy"])
    add_center_text_card(slide, 0.95, 2.75, 5.0, 1.25, "功能测试通过", "图片 / 视频 / 摄像头 / 历史记录流程完整", fill=COLORS["white"], title_size=18)
    add_center_text_card(slide, 6.75, 2.75, 5.0, 1.25, "模型可用", "主模型 YOLOv8m，mAP@0.5 达 0.829", fill=COLORS["light_orange"], title_size=18)
    add_center_text_card(slide, 0.95, 4.45, 5.0, 1.25, "性能可接受", "样例 15.17 FPS，系统平均约 9.61 FPS", fill=COLORS["light_green"], title_size=18)
    add_center_text_card(slide, 6.75, 4.45, 5.0, 1.25, "后续方向", "RGB-D、多传感器融合与轻量化部署", fill=COLORS["white"], title_size=18)
    add_bottom_note(slide, "本文工作可为果实三维定位、采摘路径规划和机械臂抓取控制提供可复用的视觉基础。")

    # Slide 12
    slide = add_base_slide(prs)
    add_title(slide, 12, "致谢")
    thank = rounded_card(slide, 0.82, 1.45, 6.0, 4.6, COLORS["white"])
    tf = thank.text_frame
    tf.clear()
    tf.margin_left = Inches(0.55)
    tf.margin_top = Inches(0.6)
    p1 = tf.paragraphs[0]
    p1.text = "谢 谢"
    set_paragraph_style(p1, size=32, bold=True, color=COLORS["navy"])
    p2 = tf.add_paragraph()
    p2.text = "感谢各位老师聆听"
    set_paragraph_style(p2, size=18, color=COLORS["muted"])
    p3 = tf.add_paragraph()
    p3.text = "恳请批评指正"
    set_paragraph_style(p3, size=18, bold=True, color=COLORS["orange"])
    add_picture_card(slide, assets["ui"], 7.05, 1.7, 5.25, 3.85, caption="系统界面示意")
    add_bottom_note(slide, "汇报结束，恳请各位老师批评指正。")

    return prs


def main() -> None:
    if not BACKUP_PATH.exists() and PPT_PATH.exists():
        shutil.copy2(PPT_PATH, BACKUP_PATH)

    assets = build_assets()
    prs = build_presentation(assets)
    prs.save(PPT_PATH)


if __name__ == "__main__":
    main()
