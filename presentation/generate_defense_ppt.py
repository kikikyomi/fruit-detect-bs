from __future__ import annotations

import json
import sqlite3
import tempfile
from collections import Counter
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DB_PATH = ROOT / "fruit-veg-detect" / "backend" / "app" / "data" / "records.sqlite3"
OUTPUT_DIR = ROOT / "fruit-veg-detect" / "backend" / "app" / "data" / "outputs"
PPT_PATH = ROOT / "presentation" / "fruit_veg_defense.pptx"
OUTLINE_PATH = ROOT / "presentation" / "defense_outline.md"


COLORS = {
    "navy": RGBColor(27, 44, 84),
    "blue": RGBColor(49, 107, 184),
    "green": RGBColor(84, 155, 98),
    "light_green": RGBColor(228, 242, 230),
    "orange": RGBColor(228, 134, 73),
    "light_orange": RGBColor(250, 235, 223),
    "bg": RGBColor(247, 249, 252),
    "text": RGBColor(46, 57, 75),
    "muted": RGBColor(108, 118, 136),
    "line": RGBColor(217, 223, 233),
    "white": RGBColor(255, 255, 255),
    "soft_blue": RGBColor(235, 242, 252),
    "soft_red": RGBColor(252, 235, 232),
}


@dataclass
class SlideSpec:
    page: int
    title: str
    bullets: list[str]
    visual: str
    remark: str


def load_metrics() -> dict[str, Any]:
    if not DB_PATH.exists():
        raise FileNotFoundError(f"Database not found: {DB_PATH}")

    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute(
        """
        SELECT id, type, created_at, file_name, input_path, output_path, result_json, summary_json
        FROM records
        ORDER BY id DESC
        """
    )
    rows = cur.fetchall()
    conn.close()

    type_counter = Counter()
    class_counter = Counter()
    tracker_counter = Counter()
    detection_counts: list[int] = []

    latest_image: dict[str, Any] | None = None
    latest_video: dict[str, Any] | None = None

    for record_id, record_type, created_at, file_name, input_path, output_path, result_json, summary_json in rows:
        type_counter[record_type] += 1
        result = json.loads(result_json) if result_json else {}
        summary = json.loads(summary_json) if summary_json else {}
        boxes = result.get("boxes") or []
        detection_counts.append(len(boxes))

        for box in boxes:
            cls_name = str(box.get("cls_name") or box.get("class_name") or box.get("label") or "unknown")
            class_counter[cls_name] += 1

        if record_type == "video":
            tracker_counter[str(summary.get("tracker", "unknown"))] += 1

        record_info = {
            "id": record_id,
            "type": record_type,
            "created_at": created_at,
            "file_name": file_name,
            "input_path": input_path,
            "output_path": output_path,
            "result": result,
            "summary": summary,
        }
        if record_type == "image" and latest_image is None:
            latest_image = record_info
        if record_type == "video" and latest_video is None:
            latest_video = record_info

    if latest_image is None or latest_video is None:
        raise RuntimeError("Need at least one image record and one video record to build the PPT.")

    latest_video_summary = latest_video["summary"]
    video_keyframe_path = None
    keyframe_details = latest_video_summary.get("keyframe_details") or []
    if keyframe_details:
        first_url = keyframe_details[0].get("image_url")
        if isinstance(first_url, str) and first_url.startswith("/static/outputs/"):
            video_keyframe_path = OUTPUT_DIR / first_url.replace("/static/outputs/", "", 1)

    return {
        "total_records": len(rows),
        "type_counter": type_counter,
        "class_counter": class_counter,
        "tracker_counter": tracker_counter,
        "avg_boxes": round(sum(detection_counts) / len(detection_counts), 1) if detection_counts else 0.0,
        "max_boxes": max(detection_counts) if detection_counts else 0,
        "latest_image": latest_image,
        "latest_video": latest_video,
        "latest_video_keyframe": video_keyframe_path,
        "snapshot_date": (rows[0][2] if rows else datetime.now().isoformat())[:10],
    }


def build_slide_specs(metrics: dict[str, Any]) -> list[SlideSpec]:
    latest_image_boxes = len(metrics["latest_image"]["result"].get("boxes") or [])
    latest_video_summary = metrics["latest_video"]["summary"]
    image_count = metrics["type_counter"].get("image", 0)
    video_count = metrics["type_counter"].get("video", 0)
    tracker_deepsort = metrics["tracker_counter"].get("deepsort", 0)
    tracker_naive = metrics["tracker_counter"].get("naive-iou", 0)

    return [
        SlideSpec(
            page=1,
            title="封面",
            bullets=[
                "题目：基于YOLO的果蔬识别系统",
                "姓名：XXX",
                "学号：XXXXXXXX",
                "导师：XXX",
                f"日期：{datetime.now():%Y年%m月%d日}",
            ],
            visual="建议放图：系统识别效果图",
            remark="开场先点题，说明这是一个面向图片、视频和摄像头场景的果蔬识别系统。",
        ),
        SlideSpec(
            page=2,
            title="目录",
            bullets=[
                "研究背景与目标",
                "系统设计与方法",
                "核心实现与效果",
                "结论与展望",
            ],
            visual="建议配图：四段式目录示意",
            remark="目录页只报结构，不展开细节，把老师注意力引到后面的系统实现和效果展示。",
        ),
        SlideSpec(
            page=3,
            title="研究背景与问题",
            bullets=[
                "人工分拣成本高",
                "果蔬外观差异大",
                "视频场景更难处理",
                "目标是做可用系统",
            ],
            visual="建议配图：痛点示意图",
            remark="先讲问题，再讲聚焦点。强调课题不是只做模型，而是把识别流程真正落地成系统。",
        ),
        SlideSpec(
            page=4,
            title="研究思路与方法",
            bullets=[
                "前后端分离架构",
                "YOLO负责目标检测",
                "DeepSORT负责跟踪",
                "SQLite保存记录",
            ],
            visual="建议放图：技术路线图",
            remark="按输入、检测、跟踪、可视化、存储五步讲，逻辑要清楚，控制在40秒内。",
        ),
        SlideSpec(
            page=5,
            title="核心研究内容1：图片识别模块",
            bullets=[
                "支持单图上传识别",
                "输出框选与类别",
                "结果自动保存入库",
                f"最新样例识别{latest_image_boxes}框",
            ],
            visual="建议放图：最新标注图片",
            remark="这一页重点看识别结果图，先说能识别，再说会输出类别、置信度和标注结果。",
        ),
        SlideSpec(
            page=6,
            title="核心研究内容2：视频与跟踪模块",
            bullets=[
                "视频按间隔抽帧",
                "支持轨迹连续绘制",
                f"样例{latest_video_summary.get('total_frames', 0)}帧用{round(float(latest_video_summary.get('processing_seconds', 0.0)), 1)}秒",
                f"处理速度{latest_video_summary.get('processing_fps', 0)}FPS",
            ],
            visual="建议放图：关键帧+速度卡片",
            remark="先报结论，系统能处理长视频并输出轨迹；再补充抽帧策略是为了平衡速度和效果。",
        ),
        SlideSpec(
            page=7,
            title="核心研究内容3：工程化与稳定性",
            bullets=[
                "记录支持回看删除",
                "默认仅留近50条",
                f"图片{image_count}条 视频{video_count}条",
                f"DeepSORT{tracker_deepsort}次 IoU{tracker_naive}次",
            ],
            visual="建议放图：记录分布图+跟踪后端图",
            remark="这一页强调系统不是演示代码，而是具备记录管理、异常回退和运行留痕能力。",
        ),
        SlideSpec(
            page=8,
            title="研究结论",
            bullets=[
                "跑通识别全流程",
                "实现检测跟踪联动",
                "具备部署与扩展性",
                "后续实验框架已留",
            ],
            visual="建议配图：结论总结卡片",
            remark="结论页只保留最强的四点，口头上再补一句：系统已具备继续做模型优化的基础。",
        ),
        SlideSpec(
            page=9,
            title="不足与展望",
            bullets=[
                "对比实验尚未固化",
                "真实采摘数据偏少",
                "类别分布不均衡",
                "后续补充新模型对比",
            ],
            visual="建议配图：问题到改进路线图",
            remark="不足一定要实事求是，但最后一句要落到可执行的改进方向，避免只停留在问题描述。",
        ),
        SlideSpec(
            page=10,
            title="致谢",
            bullets=[
                "感谢各位老师聆听",
                "恳请批评指正",
            ],
            visual="建议配图：简洁致谢页",
            remark="最后一句就够：汇报结束，恳请各位老师批评指正。",
        ),
    ]


def _set_text_style(paragraph, *, size: int, bold: bool = False, color: RGBColor | None = None, font_name: str = "Microsoft YaHei") -> None:
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = color or COLORS["text"]


def add_base_slide(prs: Presentation) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    top_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.32))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS["navy"]
    top_bar.line.fill.background()
    return slide


def add_title(slide, page: int, title: str) -> None:
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.52), Inches(0.92), Inches(0.46))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLORS["orange"]
    badge.line.fill.background()
    badge.text_frame.text = f"{page:02d}"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text_style(badge.text_frame.paragraphs[0], size=20, bold=True, color=COLORS["white"])

    title_box = slide.shapes.add_textbox(Inches(1.55), Inches(0.48), Inches(7.6), Inches(0.62))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    _set_text_style(p, size=26, bold=True, color=COLORS["navy"])


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.92), Inches(5.8), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    _set_text_style(p, size=9, color=COLORS["muted"])


def add_bullet_box(slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = COLORS["line"]
    tf = card.text_frame
    tf.clear()
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    tf.word_wrap = True
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"· {bullet}"
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        _set_text_style(p, size=22, color=COLORS["text"])


def add_remark_box(slide, remark: str) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.2), Inches(12.1), Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["soft_blue"]
    box.line.color.rgb = COLORS["line"]
    tf = box.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"备注：{remark}"
    _set_text_style(p, size=11, color=COLORS["muted"])


def add_picture_card(slide, image_path: Path, left: float, top: float, width: float, height: float, caption: str | None = None) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = COLORS["line"]
    slide.shapes.add_picture(str(image_path), Inches(left + 0.08), Inches(top + 0.08), Inches(width - 0.16), Inches(height - (0.2 if caption else 0.16)))
    if caption:
        box = slide.shapes.add_textbox(Inches(left + 0.14), Inches(top + height - 0.34), Inches(width - 0.28), Inches(0.18))
        p = box.text_frame.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        _set_text_style(p, size=9, color=COLORS["muted"])


def add_stat_card(slide, left: float, top: float, width: float, height: float, value: str, label: str, fill_color: RGBColor | None = None) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color or COLORS["white"]
    card.line.color.rgb = COLORS["line"]

    tf = card.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = value
    p1.alignment = PP_ALIGN.LEFT
    _set_text_style(p1, size=22, bold=True, color=COLORS["navy"])

    p2 = tf.add_paragraph()
    p2.text = label
    p2.alignment = PP_ALIGN.LEFT
    _set_text_style(p2, size=10, color=COLORS["muted"])


def add_flow_box(slide, left: float, top: float, width: float, height: float, title: str, subtitle: str, fill_color: RGBColor) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.fill.background()
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    _set_text_style(p1, size=17, bold=True, color=COLORS["navy"])
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.alignment = PP_ALIGN.CENTER
    _set_text_style(p2, size=10, color=COLORS["muted"])


def add_section_chip(slide, left: float, top: float, width: float, text: str, fill_color: RGBColor) -> None:
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.36))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill_color
    chip.line.fill.background()
    chip.text_frame.text = text
    chip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text_style(chip.text_frame.paragraphs[0], size=11, bold=True, color=COLORS["navy"])


def create_chart_assets(metrics: dict[str, Any], temp_dir: Path) -> dict[str, Path]:
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False

    record_chart = temp_dir / "record_mix.png"
    labels = ["图片", "视频"]
    values = [metrics["type_counter"].get("image", 0), metrics["type_counter"].get("video", 0)]
    colors = ["#4A78C2", "#E48649"]
    fig, ax = plt.subplots(figsize=(4.2, 3.0), dpi=180)
    wedges, texts, autotexts = ax.pie(
        values,
        labels=labels,
        autopct="%1.0f",
        startangle=110,
        colors=colors,
        wedgeprops={"width": 0.48, "edgecolor": "white"},
    )
    ax.set_title("运行记录分布", fontsize=12)
    for item in texts + autotexts:
        item.set_fontsize(10)
    fig.tight_layout()
    fig.savefig(record_chart, facecolor="white", bbox_inches="tight")
    plt.close(fig)

    tracker_chart = temp_dir / "tracker_mix.png"
    tracker_labels = ["DeepSORT", "naive-IoU"]
    tracker_values = [metrics["tracker_counter"].get("deepsort", 0), metrics["tracker_counter"].get("naive-iou", 0)]
    fig, ax = plt.subplots(figsize=(4.2, 3.0), dpi=180)
    bars = ax.bar(tracker_labels, tracker_values, color=["#4A78C2", "#75B07C"], width=0.55)
    ax.set_title("视频跟踪后端使用次数", fontsize=12)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="y", linestyle="--", alpha=0.25)
    ax.set_axisbelow(True)
    for bar, value in zip(bars, tracker_values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 0.3, str(value), ha="center", va="bottom", fontsize=10)
    fig.tight_layout()
    fig.savefig(tracker_chart, facecolor="white", bbox_inches="tight")
    plt.close(fig)

    image_path = Path(metrics["latest_image"]["output_path"])
    image_asset = temp_dir / "latest_image.png"
    ImageOps.exif_transpose(Image.open(image_path)).save(image_asset)

    keyframe_source = Path(metrics["latest_video_keyframe"]) if metrics["latest_video_keyframe"] else None
    if keyframe_source is None or not keyframe_source.exists():
        raise FileNotFoundError("Latest video keyframe not found.")
    keyframe_asset = temp_dir / "latest_video.png"
    keyframe = ImageOps.exif_transpose(Image.open(keyframe_source))
    if keyframe.width > keyframe.height:
        keyframe = keyframe.rotate(90, expand=True)
    keyframe.save(keyframe_asset)

    cover_asset = temp_dir / "cover.png"
    cover = ImageOps.exif_transpose(Image.open(image_path)).copy()
    cover.thumbnail((1400, 900))
    cover.save(cover_asset)

    return {
        "record_chart": record_chart,
        "tracker_chart": tracker_chart,
        "image": image_asset,
        "video": keyframe_asset,
        "cover": cover_asset,
    }


def build_presentation(metrics: dict[str, Any], specs: list[SlideSpec], assets: dict[str, Path]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 cover
    slide = add_base_slide(prs)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.0), Inches(5.5), Inches(5.6))
    hero.fill.solid()
    hero.fill.fore_color.rgb = COLORS["white"]
    hero.line.color.rgb = COLORS["line"]
    slide.shapes.add_picture(str(assets["cover"]), Inches(7.28), Inches(1.08), Inches(5.34), Inches(5.44))

    add_section_chip(slide, 0.6, 0.68, 1.25, "答辩汇报", COLORS["light_orange"])
    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.25), Inches(6.0), Inches(1.4))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "基于YOLO的果蔬识别系统设计与实现"
    _set_text_style(p1, size=28, bold=True, color=COLORS["navy"])
    p2 = tf.add_paragraph()
    p2.text = "Fruit & Vegetable Detection System"
    _set_text_style(p2, size=13, color=COLORS["muted"])

    info_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(3.0), Inches(5.45), Inches(2.0))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = COLORS["white"]
    info_box.line.color.rgb = COLORS["line"]
    info_tf = info_box.text_frame
    info_tf.margin_left = Inches(0.18)
    info_tf.margin_top = Inches(0.12)
    for idx, text in enumerate(
        [
            "姓名：XXX",
            "学号：XXXXXXXX",
            "导师：XXX",
            f"日期：{datetime.now():%Y年%m月%d日}",
        ]
    ):
        p = info_tf.paragraphs[0] if idx == 0 else info_tf.add_paragraph()
        p.text = text
        p.space_after = Pt(10)
        _set_text_style(p, size=20, color=COLORS["text"])

    add_remark_box(slide, specs[0].remark)
    add_footer(slide, "项目来源：fruit-veg-detect / 数据快照来自 records.sqlite3")

    # Slide 2 outline
    slide = add_base_slide(prs)
    add_title(slide, 2, specs[1].title)
    toc_items = specs[1].bullets
    positions = [(0.9, 1.8), (6.8, 1.8), (0.9, 3.65), (6.8, 3.65)]
    colors = [COLORS["soft_blue"], COLORS["light_orange"], COLORS["light_green"], COLORS["soft_red"]]
    for idx, ((left, top), item, color) in enumerate(zip(positions, toc_items, colors), start=1):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.4), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.text = f"{idx:02d}"
        _set_text_style(p1, size=13, bold=True, color=COLORS["orange"])
        p2 = tf.add_paragraph()
        p2.text = item
        _set_text_style(p2, size=22, bold=True, color=COLORS["navy"])
    add_remark_box(slide, specs[1].remark)

    # Slide 3 background
    slide = add_base_slide(prs)
    add_title(slide, 3, specs[2].title)
    add_bullet_box(slide, specs[2].bullets, 0.65, 1.5, 4.2, 3.7)
    pain_points = [
        ("人工环节", "效率受人工影响"),
        ("场景复杂", "遮挡与光照变化大"),
        ("实时需求", "视频更考验速度"),
    ]
    x_positions = [5.2, 8.0, 10.8]
    for (title, desc), left in zip(pain_points, x_positions):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.72), Inches(2.15), Inches(1.45))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["white"]
        card.line.color.rgb = COLORS["line"]
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.alignment = PP_ALIGN.CENTER
        _set_text_style(p1, size=16, bold=True, color=COLORS["navy"])
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.alignment = PP_ALIGN.CENTER
        _set_text_style(p2, size=11, color=COLORS["muted"])
    focus = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(3.7), Inches(7.75), Inches(1.25))
    focus.fill.solid()
    focus.fill.fore_color.rgb = COLORS["light_green"]
    focus.line.fill.background()
    tf = focus.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "研究聚焦：把识别模型接入完整系统，覆盖图片、视频、摄像头三类使用场景。"
    _set_text_style(p1, size=18, bold=True, color=COLORS["navy"])
    add_remark_box(slide, specs[2].remark)

    # Slide 4 method
    slide = add_base_slide(prs)
    add_title(slide, 4, specs[3].title)
    steps = [
        ("输入采集", "图片 / 视频 / 摄像头"),
        ("目标检测", "YOLO 推理"),
        ("目标跟踪", "DeepSORT / IoU"),
        ("结果输出", "框选 / 轨迹 / 视频"),
        ("记录管理", "SQLite 持久化"),
    ]
    fills = [COLORS["soft_blue"], COLORS["light_orange"], COLORS["light_green"], COLORS["soft_blue"], COLORS["light_orange"]]
    step_left = 0.72
    for idx, ((title, subtitle), fill) in enumerate(zip(steps, fills)):
        add_flow_box(slide, step_left + idx * 2.45, 1.78, 2.1, 1.25, title, subtitle, fill)
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(step_left + idx * 2.45 + 2.12), Inches(2.18), Inches(0.22), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["orange"]
            arrow.line.fill.background()
    add_bullet_box(slide, specs[3].bullets, 0.72, 3.55, 4.3, 2.25)
    method_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.35), Inches(3.55), Inches(7.48), Inches(2.25))
    method_card.fill.solid()
    method_card.fill.fore_color.rgb = COLORS["white"]
    method_card.line.color.rgb = COLORS["line"]
    tf = method_card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.14)
    p1 = tf.paragraphs[0]
    p1.text = "核心实现思路"
    _set_text_style(p1, size=18, bold=True, color=COLORS["navy"])
    for text in [
        "· 前端使用 Vue3 + Vite，后端使用 FastAPI。",
        "· 目标检测与 API 解耦，便于替换新权重或新模型。",
        "· 视频端加入轨迹补全与关键帧输出，增强展示效果。",
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.space_after = Pt(6)
        _set_text_style(p, size=14, color=COLORS["text"])
    add_remark_box(slide, specs[3].remark)

    # Slide 5 image module
    slide = add_base_slide(prs)
    add_title(slide, 5, specs[4].title)
    add_bullet_box(slide, specs[4].bullets, 0.68, 1.56, 4.12, 3.25)
    add_stat_card(slide, 0.68, 5.0, 1.95, 0.88, str(len(metrics["latest_image"]["result"].get("boxes") or [])), "当前样例识别框", COLORS["soft_blue"])
    add_stat_card(slide, 2.8, 5.0, 2.0, 0.88, "类别+置信度", "直接叠加可视化", COLORS["light_orange"])
    add_picture_card(slide, assets["image"], 5.12, 1.4, 7.7, 4.45, "图：最新图片识别结果")
    add_remark_box(slide, specs[4].remark)

    # Slide 6 video module
    slide = add_base_slide(prs)
    add_title(slide, 6, specs[5].title)
    add_picture_card(slide, assets["video"], 0.72, 1.45, 5.2, 4.55, "图：最新视频关键帧结果")
    video_summary = metrics["latest_video"]["summary"]
    add_stat_card(slide, 6.25, 1.58, 2.65, 0.98, str(video_summary.get("total_frames", 0)), "总帧数", COLORS["soft_blue"])
    add_stat_card(slide, 9.0, 1.58, 2.65, 0.98, str(video_summary.get("sampled_frames", 0)), "采样帧数", COLORS["light_orange"])
    add_stat_card(slide, 6.25, 2.75, 2.65, 0.98, str(video_summary.get("detections", 0)), "检测次数", COLORS["light_green"])
    add_stat_card(slide, 9.0, 2.75, 2.65, 0.98, f"{video_summary.get('processing_fps', 0)} FPS", "整体处理速度", COLORS["soft_blue"])
    add_bullet_box(slide, specs[5].bullets, 6.25, 4.02, 5.4, 1.72)
    add_footer(slide, "样例视频来源于 records.sqlite3 中最新视频记录")
    add_remark_box(slide, specs[5].remark)

    # Slide 7 engineering
    slide = add_base_slide(prs)
    add_title(slide, 7, specs[6].title)
    add_bullet_box(slide, specs[6].bullets, 0.68, 1.48, 3.95, 3.5)
    add_stat_card(slide, 0.68, 5.18, 1.85, 0.78, str(metrics["total_records"]), "当前保留记录", COLORS["light_green"])
    add_stat_card(slide, 2.65, 5.18, 1.98, 0.78, f"{metrics['avg_boxes']}", "平均框数/记录", COLORS["soft_blue"])
    add_picture_card(slide, assets["record_chart"], 5.0, 1.55, 3.7, 2.3, "图：图片/视频记录占比")
    add_picture_card(slide, assets["tracker_chart"], 8.95, 1.55, 3.7, 2.3, "图：视频跟踪后端使用次数")
    stability = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(4.2), Inches(7.65), Inches(1.4))
    stability.fill.solid()
    stability.fill.fore_color.rgb = COLORS["white"]
    stability.line.color.rgb = COLORS["line"]
    tf = stability.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.12)
    p1 = tf.paragraphs[0]
    p1.text = "稳定性设计"
    _set_text_style(p1, size=17, bold=True, color=COLORS["navy"])
    for text in [
        "· 超过保留上限后自动清理旧记录和关联文件。",
        "· 跟踪后端异常时可回退到 naive-IoU，避免流程中断。",
    ]:
        p = tf.add_paragraph()
        p.text = text
        _set_text_style(p, size=13, color=COLORS["text"])
    add_footer(slide, f"数据快照日期：{metrics['snapshot_date']}")
    add_remark_box(slide, specs[6].remark)

    # Slide 8 conclusions
    slide = add_base_slide(prs)
    add_title(slide, 8, specs[7].title)
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.35), Inches(12.0), Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLORS["soft_blue"]
    banner.line.fill.background()
    banner.text_frame.text = "结论先行：系统已经跑通“检测 - 跟踪 - 可视化 - 存储”完整闭环。"
    _set_text_style(banner.text_frame.paragraphs[0], size=22, bold=True, color=COLORS["navy"])
    cards = specs[7].bullets
    positions = [(0.95, 2.7), (6.7, 2.7), (0.95, 4.4), (6.7, 4.4)]
    fills = [COLORS["white"], COLORS["light_orange"], COLORS["light_green"], COLORS["white"]]
    for (left, top), text, fill in zip(positions, cards, fills):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.0), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        card.line.color.rgb = COLORS["line"]
        p = card.text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_text_style(p, size=20, bold=True, color=COLORS["navy"])
    add_remark_box(slide, specs[7].remark)

    # Slide 9 limitations
    slide = add_base_slide(prs)
    add_title(slide, 9, specs[8].title)
    limit_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.6), Inches(5.5), Inches(3.7))
    limit_box.fill.solid()
    limit_box.fill.fore_color.rgb = COLORS["soft_red"]
    limit_box.line.fill.background()
    tf = limit_box.text_frame
    p = tf.paragraphs[0]
    p.text = "当前不足"
    _set_text_style(p, size=22, bold=True, color=COLORS["navy"])
    for text in specs[8].bullets[:3]:
        para = tf.add_paragraph()
        para.text = f"· {text}"
        _set_text_style(para, size=18, color=COLORS["text"])

    future_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(1.6), Inches(5.5), Inches(3.7))
    future_box.fill.solid()
    future_box.fill.fore_color.rgb = COLORS["light_green"]
    future_box.line.fill.background()
    tf2 = future_box.text_frame
    p = tf2.paragraphs[0]
    p.text = "下一步改进"
    _set_text_style(p, size=22, bold=True, color=COLORS["navy"])
    for text in [
        specs[8].bullets[3],
        "补充 tomatOD / LaboroTomato 数据",
        "完善 model-compare-lab 实验结果",
    ]:
        para = tf2.add_paragraph()
        para.text = f"· {text}"
        _set_text_style(para, size=18, color=COLORS["text"])
    add_remark_box(slide, specs[8].remark)

    # Slide 10 thanks
    slide = add_base_slide(prs)
    add_title(slide, 10, specs[9].title)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.8), Inches(4.9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["white"]
    bg.line.color.rgb = COLORS["line"]
    slide.shapes.add_picture(str(assets["cover"]), Inches(7.65), Inches(1.62), Inches(4.4), Inches(4.35))
    thanks_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.05), Inches(5.9), Inches(1.8))
    tf = thanks_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "谢 谢"
    _set_text_style(p1, size=34, bold=True, color=COLORS["navy"])
    p2 = tf.add_paragraph()
    p2.text = "感谢各位老师聆听"
    _set_text_style(p2, size=18, color=COLORS["muted"])
    p3 = tf.add_paragraph()
    p3.text = "恳请批评指正"
    _set_text_style(p3, size=18, color=COLORS["orange"])
    add_remark_box(slide, specs[9].remark)

    return prs


def write_outline(specs: list[SlideSpec]) -> None:
    lines: list[str] = ["# 果蔬识别系统答辩PPT大纲", ""]
    for spec in specs:
        lines.append(f"第{spec.page}页：{spec.title}")
        lines.append("本页核心要点：")
        for bullet in spec.bullets:
            lines.append(f"· {bullet}")
        lines.append(f"建议配图/图表类型：{spec.visual}")
        lines.append(f"【备注：{spec.remark}】")
        lines.append("")
    OUTLINE_PATH.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    metrics = load_metrics()
    specs = build_slide_specs(metrics)
    write_outline(specs)

    with tempfile.TemporaryDirectory() as tmp:
        assets = create_chart_assets(metrics, Path(tmp))
        prs = build_presentation(metrics, specs, assets)
        prs.save(PPT_PATH)

    print(f"Saved PPT: {PPT_PATH}")
    print(f"Saved outline: {OUTLINE_PATH}")


if __name__ == "__main__":
    main()
